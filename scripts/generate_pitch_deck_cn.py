#!/usr/bin/env python3
"""Lumi.AI 投资者路演 PPT 生成器 — 中文版"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === COLOR PALETTE ===
BG_DARK = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
PRIMARY = RGBColor(0x63, 0x66, 0xF1)
ACCENT = RGBColor(0x8B, 0x5C, 0xF6)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MUTED = RGBColor(0x99, 0x99, 0xAA)
GREEN = RGBColor(0x10, 0xB9, 0x81)
PINK = RGBColor(0xEC, 0x48, 0x99)

W = Inches(13.333)
H = Inches(7.5)
FONT = 'Microsoft YaHei'  # Chinese font

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = FONT; p.alignment = alignment
    return txBox

def add_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_section_header(slide, number, title, subtitle=""):
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, PRIMARY)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(2), Inches(1), f"第 {number} 部分", 16, ACCENT, True)
    add_text_box(slide, Inches(1), Inches(2.8), Inches(10), Inches(1.5), title, 44, WHITE, True)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(4.3), Inches(10), Inches(1), subtitle, 20, MUTED)

# ========== SLIDE BUILDERS ==========

def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), PRIMARY)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(4), Inches(0.8), "🌟 Lumi.AI", 28, PRIMARY, True)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5), "个人命运引擎", 52, WHITE, True)
    add_text_box(slide, Inches(1), Inches(4.0), Inches(10), Inches(0.8),
                 "AI 智能键盘 × 意图交易市场 × 数字孪生优化引擎", 22, LIGHT_GRAY)
    add_text_box(slide, Inches(1), Inches(5.2), Inches(6), Inches(0.5),
                 "种子轮融资  |  $1,000,000  |  2026年2月", 16, MUTED)
    add_rect(slide, Inches(0), H - Inches(0.08), W, Inches(0.08), ACCENT)

def build_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "核心痛点", 14, ACCENT, True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(1), "AI 很强大，但不够个性化", 36, WHITE, True)
    problems = [
        ("🔒 隐私缺口", "用户将敏感数据共享给云端 AI —— 对个人信息毫无控制"),
        ("🤖 缺乏个性化", "通用 AI 回复忽略用户的价值观、风险偏好和决策风格"),
        ("📱 应用孤岛", "用户每天在 30+ 个应用间切换 —— 缺少统一的智能层"),
        ("💸 意图未被变现", "用户每天表达数百万次购买意图，但缺乏高效匹配系统"),
    ]
    for i, (title, desc) in enumerate(problems):
        left = Inches(0.8) + (i % 2) * Inches(6)
        top = Inches(2.2) + (i // 2) * Inches(2.2)
        add_rect(slide, left, top, Inches(5.5), Inches(1.8), BG_CARD, RGBColor(0x33, 0x33, 0x55))
        add_text_box(slide, left + Inches(0.3), top + Inches(0.3), Inches(4.9), Inches(0.5), title, 20, GOLD, True)
        add_text_box(slide, left + Inches(0.3), top + Inches(0.9), Inches(4.9), Inches(0.7), desc, 14, LIGHT_GRAY)

def build_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "我们的解决方案", 14, ACCENT, True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(1), "Lumi.AI —— 你的个人命运引擎", 34, WHITE, True)
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8),
                 "一款 AI 驱动的智能键盘，在输入层捕获用户意图，构建数字孪生，"
                 "驱动三层飞轮：意图层 → 超级 Agent → 意图交易市场。", 16, LIGHT_GRAY)
    layers = [
        ("第一层：智能键盘", "意图捕获 <150ms，隐私优先，本地处理，零摩擦入口", PRIMARY),
        ("第二层：数字孪生（灵魂矩阵）", "高保真用户画像 —— 人格、价值观、风险偏好、决策模式", ACCENT),
        ("第三层：超级 Agent", "贝尔曼最优决策引擎，多路径模拟，J-Curve 分析", GREEN),
        ("第四层：LIX 意图交易市场", "开放供需匹配平台 —— 商品、技能、意图，任何人都可以成为供应方", GOLD),
    ]
    for i, (title, desc, color) in enumerate(layers):
        top = Inches(2.8) + i * Inches(1.1)
        add_rect(slide, Inches(0.8), top, Inches(0.12), Inches(0.9), color)
        add_text_box(slide, Inches(1.2), top + Inches(0.05), Inches(5), Inches(0.4), title, 18, color, True)
        add_text_box(slide, Inches(1.2), top + Inches(0.45), Inches(10.5), Inches(0.45), desc, 13, LIGHT_GRAY)

# === 第一部分：技术 ===
def build_tech_section(prs):
    add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                       "01", "技术架构", "四层 Agent 架构 × 数字孪生 × 意图交易所")

def build_tech_arch_diagram(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "系统架构：端云协同超级 Agent", 30, WHITE, True)
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
                 "键盘（端侧）→ 云端编排 → 超级 Agent 执行 —— 隐私优先、低延迟设计", 14, MUTED)
    arch_img = '/tmp/arch_diagram.png'
    if os.path.exists(arch_img):
        img_w = Inches(12); img_h = Inches(12 / 3.94)
        img_left = (W - img_w) // 2; img_top = Inches(1.8)
        slide.shapes.add_picture(arch_img, img_left, img_top, img_w, img_h)
    else:
        add_text_box(slide, Inches(1), Inches(3), Inches(10), Inches(1),
                     "[架构图：keyboard_edge_cloud_superagent_arch_en.pdf]", 18, MUTED, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
                 "核心设计原则：\n"
                 "•  端侧优先：键盘哨兵 100% 本地运行（<150ms）\n"
                 "•  云端编排：超级 Agent 将任务路由到最优 AI 后端\n"
                 "•  隐私默认：原始文本永不离开设备 —— 仅传输意图向量", 13, LIGHT_GRAY)

def build_tech_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "四层 Agent 架构", 30, WHITE, True)
    layers_data = [
        ("L1：键盘哨兵", "<150ms 延迟 | 纯本地 | 不存储原文",
         "实时意图检测、PII 脱敏、Agent 触发识别", PRIMARY),
        ("L2：灵魂建筑师", "数字孪生构建器 | 渐进式学习 | 时间衰减感知",
         "构建灵魂矩阵：Big-5 人格 + 价值观 + 行为模式 + 目标", ACCENT),
        ("L3：命运推演师", "贝尔曼方程优化 | 多路径模拟",
         "期望价值计算、J-Curve 分析、后悔最小化、用户画像映射折扣因子", GREEN),
        ("L4：个人导航员", "高情商 AI | 信号灯系统（🟢🟡🛑）",
         "上下文感知语气、可操作的下一步建议、高 EQ 沟通", GOLD),
    ]
    for i, (name, specs, desc, color) in enumerate(layers_data):
        top = Inches(1.4) + i * Inches(1.45)
        add_rect(slide, Inches(0.6), top, Inches(12), Inches(1.3), BG_CARD, RGBColor(0x33, 0x33, 0x55))
        add_rect(slide, Inches(0.6), top, Inches(0.12), Inches(1.3), color)
        add_text_box(slide, Inches(1.0), top + Inches(0.1), Inches(4), Inches(0.4), name, 18, color, True)
        add_text_box(slide, Inches(6), top + Inches(0.1), Inches(6.2), Inches(0.4), specs, 12, MUTED)
        add_text_box(slide, Inches(1.0), top + Inches(0.6), Inches(11), Inches(0.5), desc, 14, LIGHT_GRAY)
    add_text_box(slide, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
                 "数据单向流动：L1 → L2 → L3 → L4  |  各层独立运行  |  隐私优先设计",
                 12, MUTED, alignment=PP_ALIGN.CENTER)

def build_tech_lix(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "LIX —— Lumi 意图交易所", 30, WHITE, True)
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.6),
                 "需求驱动市场：买家设定预算，AI Agent 广播订单，供应方选择接受或拒绝", 16, LIGHT_GRAY)
    flow_items = ["买家预算", "→", "Agent广播", "→", "供应匹配", "→", "接受/拒绝", "→", "结算"]
    x = Inches(0.4)
    for item in flow_items:
        if item == "→":
            add_text_box(slide, x, Inches(2.0), Inches(0.4), Inches(0.5), "→", 20, ACCENT, True, PP_ALIGN.CENTER)
            x += Inches(0.4)
        else:
            w = Inches(1.8) if len(item) > 4 else Inches(1.5)
            add_rect(slide, x, Inches(1.9), w, Inches(0.7), BG_CARD, PRIMARY)
            add_text_box(slide, x, Inches(2.0), w, Inches(0.5), item, 12, WHITE, False, PP_ALIGN.CENTER)
            x += w + Inches(0.1)
    metrics = [
        ("需求→匹配", "~2.5秒", "延迟"),
        ("资产类型", "2+", "实物 & 虚拟"),
        ("接受费用", "1–3%", "每笔交易"),
        ("供应方", "开放", "任何人可参与"),
    ]
    for i, (label, value, sub) in enumerate(metrics):
        left = Inches(0.8) + i * Inches(3.1)
        add_rect(slide, left, Inches(3.2), Inches(2.8), Inches(1.6), BG_CARD, RGBColor(0x33, 0x33, 0x55))
        add_text_box(slide, left, Inches(3.35), Inches(2.8), Inches(0.3), label, 12, MUTED, False, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(3.7), Inches(2.8), Inches(0.6), value, 32, GOLD, True, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(4.3), Inches(2.8), Inches(0.3), sub, 11, MUTED, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(5.2), Inches(10), Inches(0.4), "8 阶段验证管道", 18, WHITE, True)
    stages = ["来源白名单", "库存检查", "时效验证", "信任评分", "价格漂移≤30%", "超额预算", "配送可行性", "欺诈检测"]
    for i, s in enumerate(stages):
        left = Inches(0.8) + i * Inches(1.5)
        add_rect(slide, left, Inches(5.7), Inches(1.35), Inches(0.7), BG_CARD, ACCENT)
        add_text_box(slide, left, Inches(5.75), Inches(1.35), Inches(0.6), f"{i+1}. {s}", 9, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
                 "开放供应：实物商品 + 虚拟技能（咨询、设计）+ 意图（职业、投资）—— 任何人都可以成为供应方", 11, MUTED)

# === 第二部分：行业 ===
def build_industry_section(prs):
    add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                       "02", "行业分析", "肥尾经济：AGI 前夜的稀缺 AI 能力")

def build_industry_landscape(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "行业格局：AGI 前夜的窗口期", 30, WHITE, True)
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
                 "在 LIX 中，需求方拥有定价权：买家设定预算，AI Agent 广播订单，"
                 "供应方决定是否接受 —— 从根本上颠覆传统市场的定价逻辑。", 16, LIGHT_GRAY)
    headers = ["维度", "传统长尾（Amazon/Netflix）", "AI 肥尾（意图市场）"]
    for i, h in enumerate(headers):
        left = Inches(0.8) + i * Inches(4); w = Inches(3.7) if i > 0 else Inches(3.5)
        add_rect(slide, left, Inches(2.3), w, Inches(0.5), PRIMARY)
        add_text_box(slide, left, Inches(2.33), w, Inches(0.45), h, 12, WHITE, True, PP_ALIGN.CENTER)
    rows = [
        ("定价权", "供应方定价", "需求方设定预算，供应方接受/拒绝"),
        ("订单流", "买家搜索 → 挑选", "买家广播意图 → 供应方竞争"),
        ("匹配方式", "人工比价", "AI Agent 自动匹配最优报价"),
        ("稀缺性", "低（可复制商品）", "高（技能、服务、意图）"),
        ("网络效应", "双边弱", "三边强（用户-Agent-供应方）"),
        ("收入集中度", "20% 长尾 → 20% 收入", "20% 长尾 → 80% 收入"),
    ]
    for j, (dim, trad, fat) in enumerate(rows):
        top = Inches(2.85) + j * Inches(0.58)
        bg = BG_CARD if j % 2 == 0 else RGBColor(0x15, 0x15, 0x28)
        for i, val in enumerate([dim, trad, fat]):
            left = Inches(0.8) + i * Inches(4); w = Inches(3.7) if i > 0 else Inches(3.5)
            add_rect(slide, left, top, w, Inches(0.55), bg)
            color = WHITE if i == 0 else (GOLD if i == 2 else LIGHT_GRAY)
            add_text_box(slide, left + Inches(0.15), top + Inches(0.05), w - Inches(0.3), Inches(0.45), val, 11, color, i == 2)

def build_industry_demand(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "肥尾需求构成", 30, WHITE, True)
    categories = [
        ("专业咨询", "30%", "$500–5,000/次", "法律、医疗、税务规划", PRIMARY),
        ("创意生产", "25%", "$200–2,000/次", "设计、写作、技术架构", ACCENT),
        ("复杂决策", "20%", "$1,000–10,000/次", "投资、职业、教育", GREEN),
        ("个性化服务", "15%", "$100–1,000/次", "健身计划、心理咨询", GOLD),
        ("紧急任务", "10%", "2–5倍市价", "紧急翻译、危机公关", PINK),
    ]
    for i, (name, share, price, examples, color) in enumerate(categories):
        top = Inches(1.3) + i * Inches(1.15)
        add_rect(slide, Inches(0.8), top, Inches(12), Inches(1.0), BG_CARD, RGBColor(0x33, 0x33, 0x55))
        pct = int(share.replace('%', '')); bar_w = Inches(12 * pct / 100)
        r, g, b = int(str(color)[0:2], 16), int(str(color)[2:4], 16), int(str(color)[4:6], 16)
        add_rect(slide, Inches(0.8), top, bar_w, Inches(1.0), RGBColor(r // 3, g // 3, b // 3))
        add_rect(slide, Inches(0.8), top, Inches(0.1), Inches(1.0), color)
        add_text_box(slide, Inches(1.1), top + Inches(0.08), Inches(3.5), Inches(0.4), name, 16, color, True)
        add_text_box(slide, Inches(5), top + Inches(0.08), Inches(1.5), Inches(0.4), share, 22, WHITE, True)
        add_text_box(slide, Inches(7), top + Inches(0.08), Inches(2.5), Inches(0.4), price, 14, GOLD, True)
        add_text_box(slide, Inches(1.1), top + Inches(0.55), Inches(10), Inches(0.35), examples, 12, MUTED)
    add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.4),
                 "关键洞察：这些需求在 AGI 前无法被标准化 AI 完美解决 —— 必须依赖专业人类 Agent 或垂直 AI",
                 13, GOLD, True, PP_ALIGN.CENTER)

# === 第三部分：市场 ===
def build_market_section(prs):
    add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                       "03", "市场分析", "总可寻址市场 × 市场份额 × 增长轨迹")

def build_market_tam(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "总可寻址市场：$1.64 万亿", 30, WHITE, True)
    markets = [
        ("金融服务", "$20T", "25%", "15%", "$750B"),
        ("教育/培训", "$5T", "40%", "20%", "$400B"),
        ("创意/设计", "$2T", "50%", "25%", "$250B"),
        ("医疗健康", "$8T", "20%", "10%", "$160B"),
        ("法律服务", "$1T", "30%", "20%", "$60B"),
        ("企业咨询", "$300B", "35%", "20%", "$21B"),
    ]
    cols = ["领域", "全球规模", "可AI化比例", "可捕获份额", "Lumi目标市场"]
    widths = [Inches(3), Inches(2), Inches(2), Inches(2.5), Inches(2.5)]
    x = Inches(0.6)
    for c, w in zip(cols, widths):
        add_rect(slide, x, Inches(1.3), w, Inches(0.5), PRIMARY)
        add_text_box(slide, x, Inches(1.33), w, Inches(0.45), c, 12, WHITE, True, PP_ALIGN.CENTER)
        x += w
    for j, (sector, size, ai_pct, cap_pct, tam) in enumerate(markets):
        top = Inches(1.85) + j * Inches(0.58)
        bg = BG_CARD if j % 2 == 0 else RGBColor(0x15, 0x15, 0x28)
        x = Inches(0.6)
        for i, (val, w) in enumerate(zip([sector, size, ai_pct, cap_pct, tam], widths)):
            add_rect(slide, x, top, w, Inches(0.55), bg)
            color = WHITE if i == 0 else (GOLD if i == 4 else LIGHT_GRAY)
            add_text_box(slide, x + Inches(0.1), top + Inches(0.05), w - Inches(0.2), Inches(0.45),
                         val, 12 if i == 0 else 13, color, i == 4, PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
            x += w
    top = Inches(1.85) + 6 * Inches(0.58); x = Inches(0.6)
    for i, (val, w) in enumerate(zip(["合计", "$36.3T", "~30%", "~18%", "$1.64T"], widths)):
        add_rect(slide, x, top, w, Inches(0.6), RGBColor(0x2D, 0x2D, 0x5E))
        add_text_box(slide, x + Inches(0.1), top + Inches(0.08), w - Inches(0.2), Inches(0.45),
                     val, 14, GOLD if i == 4 else WHITE, True, PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
        x += w

def build_market_customers(prs):
    """目标客户群体与竞争优势"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "目标客户画像 & 为何我们能赢", 30, WHITE, True)
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
                 "谁在用 Lumi —— 以及为什么一旦开始使用就再也离不开", 15, MUTED)

    personas = [
        ("🎓 学生与研究者", "18–30", "高",
         "决策密集的学术生活：选课、职业规划、研究生申请、学术资料比价"),
        ("💼 知识工作者", "25–45", "极高",
         "每天 50+ 决策：日程安排、服务比选、供应商沟通、差旅预算"),
        ("🚀 自由职业者/创作者", "22–40", "高",
         "定价、客户匹配、项目评估 —— 收入取决于高效的意图执行"),
        ("🏢 小企业主", "30–55", "极高",
         "采购、供应商选择、预算优化 —— 每一笔支出决策都至关重要"),
        ("🛒 精明消费者", "20–50", "中高",
         "高价值购买（电子产品、旅行、保险）—— 预算驱动匹配帮用户省钱"),
    ]
    cols_h = ["人群", "年龄", "意图密度", "Lumi 适配原因"]
    widths = [Inches(2.8), Inches(1), Inches(1.5), Inches(7)]
    x = Inches(0.5)
    for c, w in zip(cols_h, widths):
        add_rect(slide, x, Inches(1.6), w, Inches(0.45), PRIMARY)
        add_text_box(slide, x, Inches(1.62), w, Inches(0.4), c, 11, WHITE, True, PP_ALIGN.CENTER)
        x += w
    for j, (seg, age, density, reason) in enumerate(personas):
        top = Inches(2.1) + j * Inches(0.58)
        bg = BG_CARD if j % 2 == 0 else RGBColor(0x15, 0x15, 0x28)
        x = Inches(0.5)
        for i, (val, w) in enumerate(zip([seg, age, density, reason], widths)):
            add_rect(slide, x, top, w, Inches(0.55), bg)
            color = WHITE if i == 0 else (GOLD if i == 2 else LIGHT_GRAY)
            add_text_box(slide, x + Inches(0.1), top + Inches(0.08), w - Inches(0.2), Inches(0.4),
                         val, 11, color, i == 0, PP_ALIGN.LEFT if i in (0, 3) else PP_ALIGN.CENTER)
            x += w

    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(6), Inches(0.4),
                 "客户群体共性特征", 18, WHITE, True)
    traits = [
        "•  移动优先：每天 4+ 小时键盘使用",
        "•  决策密集：每天 30–80 次意图行为",
        "•  在乎价格但缺少时间：愿意为更好结果付费",
        "•  注重隐私：偏好端侧 AI，不愿数据上云",
    ]
    for i, t in enumerate(traits):
        add_text_box(slide, Inches(0.7), Inches(5.5) + i * Inches(0.35), Inches(5.5), Inches(0.3),
                     t, 11, LIGHT_GRAY)

    add_rect(slide, Inches(6.8), Inches(5.1), Inches(5.8), Inches(2.2), BG_CARD, GOLD)
    add_text_box(slide, Inches(7.1), Inches(5.2), Inches(5.2), Inches(0.4),
                 "🏆 我们为什么能赢", 16, GOLD, True)
    reasons = [
        "•  零摩擦入口：键盘 = 通用输入层，无需下载新 App",
        "•  数字孪生锁定：迁移成本随使用时间增长",
        "•  预算驱动匹配：用户立即享受省钱体验",
        "•  网络效应：用户越多 = 供应方竞争越激烈 = 价格更优",
        "•  习惯养成：每天键盘使用 = 每天活跃参与",
    ]
    for i, r in enumerate(reasons):
        add_text_box(slide, Inches(7.3), Inches(5.65) + i * Inches(0.32), Inches(5.1), Inches(0.3),
                     r, 11, LIGHT_GRAY)

def build_market_why_now(prs):
    """为什么是现在 — OpenClaw 已验证 PMF"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "为什么是现在", 14, ACCENT, True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(1),
                 "OpenClaw 已验证 PMF — 直接进入增长期", 32, WHITE, True)

    # 左侧卡片：OpenClaw PMF 证据
    add_rect(slide, Inches(0.6), Inches(1.9), Inches(5.8), Inches(4.8), BG_CARD, GREEN)
    add_text_box(slide, Inches(0.9), Inches(2.05), Inches(5.2), Inches(0.5),
                 "✅ OpenClaw：PMF 已验证", 20, GREEN, True)
    add_text_box(slide, Inches(0.9), Inches(2.6), Inches(5.2), Inches(0.5),
                 "OpenClaw 已证明全新的 AI-Agent 工作逻辑切实可行：", 13, LIGHT_GRAY)

    evidence = [
        ("全新工作逻辑", "意图驱动 → Agent 执行 → 供需匹配 —— 完整链路已跑通"),
        ("用户行为验证", "用户自发采用 Agent 工作流，替代传统手动搜索和比价"),
        ("留存信号强劲", "高复访率 —— 用户因省时省钱而持续使用，形成习惯"),
        ("市场自发拉动", "无需付费获客即实现有机增长 —— 需求拉动，而非推送"),
    ]
    for i, (label, desc) in enumerate(evidence):
        top = Inches(3.2) + i * Inches(0.75)
        add_rect(slide, Inches(0.9), top, Inches(0.1), Inches(0.6), GREEN)
        add_text_box(slide, Inches(1.2), top + Inches(0.02), Inches(4.8), Inches(0.3),
                     label, 14, GOLD, True)
        add_text_box(slide, Inches(1.2), top + Inches(0.32), Inches(4.8), Inches(0.3),
                     desc, 11, LIGHT_GRAY)

    # 核心结论
    add_rect(slide, Inches(0.9), Inches(6.25), Inches(5.2), Inches(0.35), RGBColor(0x10, 0x3D, 0x2D))
    add_text_box(slide, Inches(1.0), Inches(6.27), Inches(5.0), Inches(0.3),
                 "\"市场存在，逻辑成立，用户留存。\"", 13, GREEN, True, PP_ALIGN.CENTER)

    # 右侧卡片：为什么现在进入增长期
    add_rect(slide, Inches(6.8), Inches(1.9), Inches(5.8), Inches(4.8), BG_CARD, GOLD)
    add_text_box(slide, Inches(7.1), Inches(2.05), Inches(5.2), Inches(0.5),
                 "🚀 跳过验证 → 直接进入增长", 20, GOLD, True)
    add_text_box(slide, Inches(7.1), Inches(2.6), Inches(5.2), Inches(0.5),
                 "Lumi 无需再证明市场 —— 市场已经被证明：", 13, LIGHT_GRAY)

    reasons = [
        ("✅ PMF 已验证", "OpenClaw 证明意图驱动的 Agent 工作流有真实需求 —— 无需再探索"),
        ("✅ 用户行为已形成", "全新工作逻辑（AI 执行，用户定目标）已被接受并形成习惯"),
        ("✅ 时机窗口最佳", "AI Agent 浪潮正处高峰 —— 抢占心智的窗口是现在，不是以后"),
        ("✅ 直接规模化", "逻辑已验证，每一分钱都投入增长，而非试错"),
    ]
    for i, (title, desc) in enumerate(reasons):
        top = Inches(3.2) + i * Inches(0.75)
        add_text_box(slide, Inches(7.3), top + Inches(0.02), Inches(5.0), Inches(0.3),
                     title, 14, WHITE, True)
        add_text_box(slide, Inches(7.3), top + Inches(0.32), Inches(5.0), Inches(0.3),
                     desc, 11, LIGHT_GRAY)

    # 底部横幅
    add_rect(slide, Inches(0.6), Inches(6.9), Inches(12), Inches(0.45), RGBColor(0x2D, 0x2D, 0x5E))
    add_text_box(slide, Inches(0.6), Inches(6.93), Inches(12), Inches(0.4),
                 "OpenClaw 验证了逻辑。 Lumi 扩展网络。 市场份额由我们来拿。",
                 15, GOLD, True, PP_ALIGN.CENTER)


def build_market_share(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "市场份额轨迹（心智驱动的非线性增长）", 28, WHITE, True)
    years = [
        ("Y1", "50亿", "2%", "1%", "100万", "早期用户，开发者社区"),
        ("Y2", "52亿", "5%", "2%", "500万", "效率工具博主，口碑传播"),
        ("Y3*", "54亿", "10%", "4%", "2000万", "心智临界点 —— 大规模采用开始"),
        ("Y4", "56亿", "15%", "6%", "5000万", "主流增长"),
        ("Y5", "58亿", "20%", "8%", "1亿", "平台级规模"),
        ("Y7", "62亿", "30%", "10%", "2亿", "意图经济基础设施"),
    ]
    cols_h = ["年份", "全球用户", "Agent OS渗透率", "Lumi份额", "Lumi用户数", "阶段"]
    widths = [Inches(1), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(4)]
    x = Inches(0.5)
    for c, w in zip(cols_h, widths):
        add_rect(slide, x, Inches(1.3), w, Inches(0.5), PRIMARY)
        add_text_box(slide, x, Inches(1.33), w, Inches(0.45), c, 11, WHITE, True, PP_ALIGN.CENTER)
        x += w
    for j, row in enumerate(years):
        top = Inches(1.85) + j * Inches(0.65)
        bg = RGBColor(0x1E, 0x1E, 0x40) if row[0] == "Y3*" else (BG_CARD if j % 2 == 0 else RGBColor(0x15, 0x15, 0x28))
        x = Inches(0.5)
        for i, (val, w) in enumerate(zip(row, widths)):
            add_rect(slide, x, top, w, Inches(0.6), bg)
            is_hl = row[0] == "Y3*"
            color = GOLD if is_hl else (WHITE if i in (0, 4) else LIGHT_GRAY)
            add_text_box(slide, x + Inches(0.05), top + Inches(0.08), w - Inches(0.1), Inches(0.45),
                         val, 11, color, is_hl or i == 4, PP_ALIGN.CENTER if i < 5 else PP_ALIGN.LEFT)
            x += w
    add_text_box(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(0.5),
                 "* Y3 = 心智临界点 —— 对标抖音 2019 年拐点，大规模采用开始",
                 14, GOLD, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.8),
                 "护城河：数字孪生锁定（迁移成本）+ 习惯路径依赖 + Agent 市场生态 + 先发品牌（\"Agent = Lumi\"）",
                 13, MUTED, False, PP_ALIGN.CENTER)

# === 第四部分：融资计划 ===
def build_financing_section(prs):
    add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                       "04", "融资计划", "种子轮：$1,000,000 出让 10% 股权  |  估值：$10,000,000")

def build_business_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "商业模式：双引擎收入", 30, WHITE, True)
    add_rect(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(2.8), BG_CARD, PRIMARY)
    add_text_box(slide, Inches(0.9), Inches(1.45), Inches(5), Inches(0.5), "💎 收入流 1：键盘订阅", 20, PRIMARY, True)
    for i, item in enumerate([
        "免费增值模式：基础免费 → 高级 AI 功能付费",
        "高级版：$9.99–19.99/月/用户",
        "功能：数字孪生、命运引擎、优先 Agent",
        "目标：20–40% 付费转化率",
        "经常性 SaaS 收入，高留存",
    ]):
        add_text_box(slide, Inches(1.1), Inches(2.1) + i * Inches(0.38), Inches(5), Inches(0.35), f"•  {item}", 12, LIGHT_GRAY)
    add_rect(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.8), BG_CARD, GOLD)
    add_text_box(slide, Inches(7.1), Inches(1.45), Inches(5), Inches(0.5), "🏪 收入流 2：LIX 交易费", 20, GOLD, True)
    for i, item in enumerate([
        "1%–3% 撮合成交费",
        "任何人都可以成为供应方 —— 开放市场",
        "收入随 GMV 增长（网络效应）",
        "8 阶段验证保证质量 → 信任 → 交易量",
        "商品 + 技能 + 服务 + 意图",
    ]):
        add_text_box(slide, Inches(7.3), Inches(2.1) + i * Inches(0.38), Inches(5), Inches(0.35), f"•  {item}", 12, LIGHT_GRAY)
    add_text_box(slide, Inches(0.8), Inches(4.4), Inches(10), Inches(0.5), "收入预测", 22, WHITE, True)
    rev_cols = ["年份", "订阅收入", "LIX GMV", "平台抽成（1-3%）", "总收入"]
    rev_widths = [Inches(1.5), Inches(2.5), Inches(2.5), Inches(3), Inches(2.5)]
    x = Inches(0.6)
    for c, w in zip(rev_cols, rev_widths):
        add_rect(slide, x, Inches(4.9), w, Inches(0.45), PRIMARY)
        add_text_box(slide, x, Inches(4.92), w, Inches(0.4), c, 11, WHITE, True, PP_ALIGN.CENTER)
        x += w
    for j, row in enumerate([("Y1","$12M","$2M","$60K","$12M"),("Y3","$300M","$500M","$15M","$315M"),("Y5","$2B","$10B","$300M","$2.3B")]):
        top = Inches(5.4) + j * Inches(0.5)
        bg = BG_CARD if j % 2 == 0 else RGBColor(0x15, 0x15, 0x28); x = Inches(0.6)
        for i, (v, w) in enumerate(zip(row, rev_widths)):
            add_rect(slide, x, top, w, Inches(0.48), bg)
            color = GOLD if i == 4 else (WHITE if i == 0 else LIGHT_GRAY)
            add_text_box(slide, x, top + Inches(0.03), w, Inches(0.4), v, 12, color, i == 4, PP_ALIGN.CENTER)
            x += w
    add_text_box(slide, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
                 "关键拐点：Y5 后，LIX 市场收入超越订阅收入 —— 平台经济效应显现", 12, GOLD, True, PP_ALIGN.CENTER)

def build_use_of_funds(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "资金用途与投资回报", 30, WHITE, True)
    add_rect(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(2.5), BG_CARD, PRIMARY)
    add_text_box(slide, Inches(0.9), Inches(1.45), Inches(5), Inches(0.5), "种子轮条款", 20, PRIMARY, True)
    for i, (label, val) in enumerate([("融资金额","$1,000,000"),("出让股权","10%"),("投前估值","$10,000,000"),("投后估值","$11,000,000")]):
        add_text_box(slide, Inches(1.1), Inches(2.1)+i*Inches(0.45), Inches(3), Inches(0.4), label, 14, LIGHT_GRAY)
        add_text_box(slide, Inches(4.2), Inches(2.1)+i*Inches(0.45), Inches(2), Inches(0.4), val, 14, GOLD, True)
    # 资金分配 — 详细明细
    add_rect(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.8), BG_CARD, ACCENT)
    add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5), Inches(0.4),
                 "资金分配 — 详细明细", 18, ACCENT, True)
    allocs = [
        ("办公场地", "$50,000", "5%", LIGHT_GRAY),
        ("设备与基础设施", "$100,000", "10%", LIGHT_GRAY),
        ("端侧模型训练", "$100,000", "10%", PRIMARY),
        ("CTO / 技术负责人", "$100,000", "10%", PRIMARY),
        ("法律合规", "$100,000", "10%", LIGHT_GRAY),
        ("市场推广与获客", "$200,000", "20%", GOLD),
        ("流动资金（储备）", "$350,000", "35%", GREEN),
    ]
    for i, (label, amt, pct, color) in enumerate(allocs):
        top = Inches(1.85) + i * Inches(0.3)
        add_text_box(slide, Inches(7.2), top, Inches(2.8), Inches(0.28), label, 10, color, False)
        add_text_box(slide, Inches(10.2), top, Inches(1.1), Inches(0.28), amt, 10, WHITE, True)
        add_text_box(slide, Inches(11.5), top, Inches(0.8), Inches(0.28), pct, 10, GOLD, True)
    # 合计行
    add_rect(slide, Inches(7.1), Inches(3.98), Inches(5.3), Inches(0.02), ACCENT)
    add_text_box(slide, Inches(7.2), Inches(4.0), Inches(2.8), Inches(0.28), "合计", 11, WHITE, True)
    add_text_box(slide, Inches(10.2), Inches(4.0), Inches(1.1), Inches(0.28), "$1,000,000", 11, GOLD, True)
    add_text_box(slide, Inches(11.5), Inches(4.0), Inches(0.8), Inches(0.28), "100%", 11, GOLD, True)
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.5), "投资回报情景分析", 22, WHITE, True)
    scenarios = [
        ("保守情景","30%","小众工具（<5000万用户）","3×","10×",LIGHT_GRAY),
        ("基准情景","45%","主流平台（1-5亿用户）","10×","50×",WHITE),
        ("乐观情景","25%","基础设施（5亿+用户）","50×","200×",GOLD),
    ]
    ws = [Inches(1.8),Inches(1.5),Inches(4),Inches(1.5),Inches(1.5)]
    x = Inches(0.8)
    for c, w in zip(["情景","概率","结果","回报","规模化"], ws):
        add_rect(slide, x, Inches(4.8), w, Inches(0.45), PRIMARY)
        add_text_box(slide, x, Inches(4.82), w, Inches(0.4), c, 11, WHITE, True, PP_ALIGN.CENTER)
        x += w
    for j, (scen, prob, outcome, ret, scale, color) in enumerate(scenarios):
        top = Inches(5.3)+j*Inches(0.55)
        bg = BG_CARD if j % 2 == 0 else RGBColor(0x15,0x15,0x28); x = Inches(0.8)
        for i, (v, w) in enumerate(zip([scen,prob,outcome,ret,scale], ws)):
            add_rect(slide, x, top, w, Inches(0.5), bg)
            c = GOLD if i >= 3 else (color if i == 0 else LIGHT_GRAY)
            add_text_box(slide, x+Inches(0.05), top+Inches(0.05), w-Inches(0.1), Inches(0.4),
                         v, 11, c, i >= 3, PP_ALIGN.CENTER if i != 2 else PP_ALIGN.LEFT)
            x += w
    add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.4),
                 "期望回报（加权）：15.4× —— 基于风险调整后的概率分析", 14, GOLD, True, PP_ALIGN.CENTER)

def build_milestones(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "关键里程碑与企业管线", 30, WHITE, True)
    milestones = [
        ("2024年", "伦敦成立公司", "公司注册，初始研发"),
        ("2025.12", "产品概念验证", "Web 模拟器 + Android 键盘原型"),
        ("2026.02", "产品演示就绪", "完整 LIX MVP、四层 Agent、DTOE v0.3"),
        ("2026 Q2", "种子轮关闭", "目标 100 万用户，LIX 市场上线"),
        ("2026 Q4", "增长阶段", "500 万用户，扩展供应网络"),
        ("2027年", "A 轮就绪", "2000 万用户，心智临界点"),
    ]
    for i, (date, event, desc) in enumerate(milestones):
        top = Inches(1.3)+i*Inches(0.95)
        add_rect(slide, Inches(2.5), top, Inches(0.04), Inches(0.95), RGBColor(0x33,0x33,0x55))
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.38), top+Inches(0.2), Inches(0.28), Inches(0.28))
        dot.fill.solid(); dot.fill.fore_color.rgb = PRIMARY; dot.line.fill.background()
        add_text_box(slide, Inches(0.6), top+Inches(0.15), Inches(1.7), Inches(0.4), date, 14, GOLD, True, PP_ALIGN.RIGHT)
        add_text_box(slide, Inches(3.0), top+Inches(0.1), Inches(4), Inches(0.35), event, 16, WHITE, True)
        add_text_box(slide, Inches(3.0), top+Inches(0.48), Inches(9), Inches(0.35), desc, 12, MUTED)
    add_rect(slide, Inches(8), Inches(1.5), Inches(4.5), Inches(2.5), BG_CARD, GOLD)
    add_text_box(slide, Inches(8.3), Inches(1.65), Inches(4), Inches(0.4), "🏦 企业管线", 16, GOLD, True)
    add_text_box(slide, Inches(8.3), Inches(2.1), Inches(4), Inches(1.6),
                 "• 摩根士丹利 — 银行级 AI 执行层\n• 90 天 PoC 框架已定义\n• 财富管理与研究工作流\n• 数字孪生优化顾问服务", 12, LIGHT_GRAY)

# === 第五部分：团队 ===
def build_team_section(prs):
    add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                       "05", "团队介绍", "UCL 创新 × 风险投资专业背景")

def build_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "创始团队", 30, WHITE, True)
    # 李松益
    add_rect(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(4.5), BG_CARD, GOLD)
    a1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.6), Inches(1.6), Inches(1.2), Inches(1.2))
    a1.fill.solid(); a1.fill.fore_color.rgb = GOLD; a1.line.fill.background()
    add_text_box(slide, Inches(1.6), Inches(1.85), Inches(1.2), Inches(0.7), "SL", 28, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.1), Inches(1.6), Inches(3), Inches(0.4), "李松益", 24, WHITE, True)
    add_text_box(slide, Inches(3.1), Inches(2.1), Inches(3), Inches(0.35), "创始人 & CEO", 16, GOLD, True)
    for i, e in enumerate(["🎓 博士在读，金融科技 — UCL", "🎓 硕士，风险投资与私募股权 — UCL", "🎓 学士，数学 — 西交利物浦 & 利物浦大学"]):
        add_text_box(slide, Inches(1.0), Inches(3.1)+i*Inches(0.4), Inches(5), Inches(0.35), e, 12, LIGHT_GRAY)
    add_text_box(slide, Inches(1.0), Inches(4.5), Inches(5), Inches(0.6),
                 "引领金融科技创新与 AI 生产力工具的融合。在金融科技研究和风险投资领域拥有深厚专业背景。", 11, MUTED)
    # 赵博
    add_rect(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(4.5), BG_CARD, ACCENT)
    a2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.8), Inches(1.6), Inches(1.2), Inches(1.2))
    a2.fill.solid(); a2.fill.fore_color.rgb = ACCENT; a2.line.fill.background()
    add_text_box(slide, Inches(7.8), Inches(1.85), Inches(1.2), Inches(0.7), "BZ", 28, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(9.3), Inches(1.6), Inches(3), Inches(0.4), "赵博", 24, WHITE, True)
    add_text_box(slide, Inches(9.3), Inches(2.1), Inches(3), Inches(0.35), "联合创始人 & COO", 16, ACCENT, True)
    for i, e in enumerate(["🎓 硕士，风险投资与私募股权 — UCL", "🎓 学士，金融学 — 杜伦大学", "💼 前国信证券金融专业人士"]):
        add_text_box(slide, Inches(7.2), Inches(3.1)+i*Inches(0.4), Inches(5), Inches(0.35), e, 12, LIGHT_GRAY)
    add_text_box(slide, Inches(7.2), Inches(4.5), Inches(5), Inches(0.6),
                 "将风险投资专业知识和金融洞察力应用于 Lumi.AI 的全球化扩张。在融资和财务运营方面经验丰富。", 11, MUTED)
    add_rect(slide, Inches(0.6), Inches(6.1), Inches(12), Inches(1.0), BG_CARD, RGBColor(0x33,0x33,0x55))
    add_text_box(slide, Inches(1.0), Inches(6.2), Inches(5), Inches(0.4), "🏛️  学术合作伙伴：伦敦大学学院（UCL）", 16, WHITE, True)
    add_text_box(slide, Inches(1.0), Inches(6.6), Inches(11), Inches(0.35),
                 "全球 Top 10 大学  |  金融科技研究中心  |  英国伦敦  |  聚焦：FinTech × AI × 隐私科技", 12, MUTED)

def build_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_rect(slide, Inches(0), H-Inches(0.06), W, Inches(0.06), ACCENT)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1), "🌟 Lumi.AI", 36, PRIMARY, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1), "构建意图经济的基础设施", 32, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1),
                 "\"Lumi 不是在卖工具 —— 是在构建 AGI 前夜的稀缺能力分配网络。\"", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    for i, (num, label) in enumerate([("$10M","估值"),("$1M","融资金额"),("10%","出让股权"),("15.4×","期望回报")]):
        left = Inches(1.5)+i*Inches(2.8)
        add_rect(slide, left, Inches(5.0), Inches(2.3), Inches(1.2), BG_CARD, PRIMARY)
        add_text_box(slide, left, Inches(5.1), Inches(2.3), Inches(0.6), num, 28, GOLD, True, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(5.7), Inches(2.3), Inches(0.4), label, 13, MUTED, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(6.6), Inches(9), Inches(0.5),
                 "联系方式：team@lumi.ai  |  lumi-agent-simulator.vercel.app  |  英国伦敦", 14, MUTED, False, PP_ALIGN.CENTER)

# ========== MAIN ==========
def main():
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H

    build_cover(prs); build_problem(prs); build_solution(prs)

    build_tech_section(prs); build_tech_arch_diagram(prs); build_tech_arch(prs); build_tech_lix(prs)

    build_industry_section(prs); build_industry_landscape(prs); build_industry_demand(prs)

    build_market_section(prs); build_market_tam(prs); build_market_customers(prs); build_market_why_now(prs); build_market_share(prs)

    build_financing_section(prs); build_business_model(prs); build_use_of_funds(prs)

    build_team_section(prs); build_team(prs)

    build_closing(prs)

    output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)),
                               "Lumi_AI_投资者路演_中文版.pptx")
    prs.save(output_path)
    print(f"✅ 中文版路演 PPT 已保存至：{output_path}")
    print(f"   总幻灯片数：{len(prs.slides)}")

if __name__ == "__main__":
    main()
