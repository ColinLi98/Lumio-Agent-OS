#!/usr/bin/env python3
"""Lumi.AI Investor Pitch Deck Generator - Part 1: Helpers & Config"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === COLOR PALETTE ===
BG_DARK = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
PRIMARY = RGBColor(0x63, 0x66, 0xF1)
ACCENT = RGBColor(0x8B, 0x5C, 0xF6)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MUTED = RGBColor(0x99, 0x99, 0xAA)
GREEN = RGBColor(0x10, 0xB9, 0x81)
PINK = RGBColor(0xEC, 0x48, 0x99)

W = Inches(13.333)  # 16:9 widescreen
H = Inches(7.5)

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(text_frame, text, font_size=14, color=LIGHT_GRAY, bold=False,
             alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(2)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_section_header(slide, number, title, subtitle=""):
    set_slide_bg(slide)
    # Accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, PRIMARY)
    # Section number
    add_text_box(slide, Inches(1), Inches(2.2), Inches(2), Inches(1),
                 f"PART {number}", 16, ACCENT, True)
    # Title
    add_text_box(slide, Inches(1), Inches(2.8), Inches(10), Inches(1.5),
                 title, 44, WHITE, True)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(4.3), Inches(10), Inches(1),
                     subtitle, 20, MUTED)

# ========== SLIDE BUILDERS ==========

def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    # Decorative gradient bar at top
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), PRIMARY)
    # Logo area
    add_text_box(slide, Inches(1), Inches(1.5), Inches(4), Inches(0.8),
                 "🌟 Lumi.AI", 28, PRIMARY, True)
    # Main title
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                 "Personal Destiny Engine", 52, WHITE, True)
    # Subtitle
    add_text_box(slide, Inches(1), Inches(4.0), Inches(10), Inches(0.8),
                 "AI-Powered Smart Keyboard × Intent Exchange Market × Digital Twin Optimization",
                 22, LIGHT_GRAY)
    # Details
    add_text_box(slide, Inches(1), Inches(5.2), Inches(6), Inches(0.5),
                 "Seed Round  |  $1,000,000  |  February 2026", 16, MUTED)
    # Bottom bar
    add_rect(slide, Inches(0), H - Inches(0.08), W, Inches(0.08), ACCENT)


def build_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 "THE PROBLEM", 14, ACCENT, True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(1),
                 "AI is Powerful, But Not Personal", 36, WHITE, True)

    problems = [
        ("🔒 Privacy Gap", "Users share sensitive data with cloud AI — no control over personal information"),
        ("🤖 No Personalization", "Generic AI responses ignore user's values, risk tolerance, and decision style"),
        ("📱 App Silos", "Users juggle 30+ apps daily — no unified intelligent layer across all interactions"),
        ("💸 No Intent Monetization", "Users express purchase intent millions of times daily, but no efficient matching system exists"),
    ]
    for i, (title, desc) in enumerate(problems):
        left = Inches(0.8) + (i % 2) * Inches(6)
        top = Inches(2.2) + (i // 2) * Inches(2.2)
        card = add_rect(slide, left, top, Inches(5.5), Inches(1.8), BG_CARD, RGBColor(0x33, 0x33, 0x55))
        tb = add_text_box(slide, left + Inches(0.3), top + Inches(0.3), Inches(4.9), Inches(0.5),
                          title, 20, GOLD, True)
        add_text_box(slide, left + Inches(0.3), top + Inches(0.9), Inches(4.9), Inches(0.7),
                     desc, 14, LIGHT_GRAY)


def build_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 "OUR SOLUTION", 14, ACCENT, True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(1),
                 "Lumi.AI — Your Personal Destiny Engine", 34, WHITE, True)

    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8),
                 "An AI-powered smart keyboard that captures user intent at the input layer, builds a Digital Twin, "
                 "and powers a three-layer flywheel: Intent Layer → Super Agent → Intent Exchange Market.",
                 16, LIGHT_GRAY)

    layers = [
        ("Layer 1: Smart Keyboard", "Intent capture at <150ms, privacy-first, local processing, zero-friction entry point", PRIMARY),
        ("Layer 2: Digital Twin (Soul Matrix)", "High-fidelity user profile — personality, values, risk tolerance, decision patterns", ACCENT),
        ("Layer 3: Super Agent", "Bellman-optimized decision engine, multi-path simulation, J-Curve analysis", GREEN),
        ("Layer 4: LIX Intent Exchange", "Open demand-supply matching for goods, skills & intentions — anyone can be a supplier", GOLD),
    ]
    for i, (title, desc, color) in enumerate(layers):
        top = Inches(2.8) + i * Inches(1.1)
        add_rect(slide, Inches(0.8), top, Inches(0.12), Inches(0.9), color)
        add_text_box(slide, Inches(1.2), top + Inches(0.05), Inches(5), Inches(0.4),
                     title, 18, color, True)
        add_text_box(slide, Inches(1.2), top + Inches(0.45), Inches(10.5), Inches(0.45),
                     desc, 13, LIGHT_GRAY)


# === PART 1: TECHNOLOGY ===
def build_tech_section(prs):
    add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                       "01", "TECHNOLOGY",
                       "Four-Layer Agent Architecture × Digital Twin × Intent Exchange")

def build_tech_arch_diagram(prs):
    """Slide with the Edge-Cloud Super Agent Architecture diagram from PDF."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "System Architecture: Edge-Cloud Super Agent", 30, WHITE, True)
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
                 "Keyboard (Edge) → Cloud Orchestration → Super Agent Execution — privacy-first, low-latency design",
                 14, MUTED)
    # Embed the architecture diagram image
    arch_img = '/tmp/arch_diagram.png'
    if os.path.exists(arch_img):
        # Image is 6444x1635 (ratio ~3.94:1), fit within slide width with padding
        img_w = Inches(12)
        img_h = Inches(12 / 3.94)  # maintain aspect ratio ~3.05"
        img_left = (W - img_w) // 2
        img_top = Inches(1.8)
        slide.shapes.add_picture(arch_img, img_left, img_top, img_w, img_h)
    else:
        add_text_box(slide, Inches(1), Inches(3), Inches(10), Inches(1),
                     "[Architecture diagram: keyboard_edge_cloud_superagent_arch_en.pdf]",
                     18, MUTED, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
                 "Key Design Principles:\n"
                 "•  Edge-first processing: Keyboard Sentinel runs 100% locally (<150ms)\n"
                 "•  Cloud orchestration: Super Agent routes tasks to optimal AI backends\n"
                 "•  Privacy-by-default: Raw text never leaves device — only intent vectors transmitted",
                 13, LIGHT_GRAY)

def build_tech_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Four-Layer Agent Architecture", 30, WHITE, True)

    layers_data = [
        ("L1: Keyboard Sentinel", "<150ms latency | Local-only | No raw text stored",
         "Real-time intent detection, PII masking, agent trigger recognition", PRIMARY),
        ("L2: Soul Architect", "Digital Twin builder | Progressive learning | Time-decay aware",
         "Builds SoulMatrix: Big-5 personality + values + behavioral patterns + goals", ACCENT),
        ("L3: Destiny Engine", "Bellman equation optimization | Multi-path simulation",
         "EV calculation, J-Curve analysis, regret minimization, gamma-mapped to user profile", GREEN),
        ("L4: Personal Navigator", "Empathetic AI | Traffic-light system (🟢🟡🛑)",
         "Context-aware tone, actionable next steps, high-EQ communication", GOLD),
    ]
    for i, (name, specs, desc, color) in enumerate(layers_data):
        top = Inches(1.4) + i * Inches(1.45)
        card = add_rect(slide, Inches(0.6), top, Inches(12), Inches(1.3), BG_CARD, RGBColor(0x33, 0x33, 0x55))
        # Color indicator
        add_rect(slide, Inches(0.6), top, Inches(0.12), Inches(1.3), color)
        add_text_box(slide, Inches(1.0), top + Inches(0.1), Inches(4), Inches(0.4),
                     name, 18, color, True)
        add_text_box(slide, Inches(6), top + Inches(0.1), Inches(6.2), Inches(0.4),
                     specs, 12, MUTED)
        add_text_box(slide, Inches(1.0), top + Inches(0.6), Inches(11), Inches(0.5),
                     desc, 14, LIGHT_GRAY)

    add_text_box(slide, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
                 "Data flows unidirectionally: L1 → L2 → L3 → L4  |  Each layer operates independently  |  Privacy-by-default",
                 12, MUTED, alignment=PP_ALIGN.CENTER)


def build_tech_digital_twin(prs):
    """Digital Twin deep-dive: 13-dim state vector, quantification, particle filter, LIX bridge."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)
    add_text_box(slide, Inches(0.8), Inches(0.35), Inches(10), Inches(0.7),
                 "Digital Twin — Soul Matrix Architecture", 30, WHITE, True)
    add_text_box(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.5),
                 "13-Dimension Bayesian Belief State · Particle Filter Inference · Privacy-Preserving Personalization",
                 14, MUTED)

    # --- Left panel: State Vector table ---
    dims = [
        ("\U0001f4b0 Wealth",       "Budget data, spending patterns",          "0.65"),
        ("\u2764\ufe0f Health",     "Health queries, fitness app usage",       "0.72"),
        ("\U0001f9e0 Skill",        "Learning intent frequency",               "0.80"),
        ("\u26a1 Energy",           "Typing speed, session timing",            "0.55"),
        ("\U0001f91d Social",       "Social intent frequency",                 "0.58"),
        ("\U0001f4bc Career",       "Job-related queries",                     "0.71"),
        ("\u2b50 Reputation",       "Social signal quality",                   "0.63"),
        ("\U0001f513 Optionality",   "Intent domain diversity",                 "0.48"),
        ("\U0001f630 Stress",       "Negative words, late-night use",          "0.62"),
        ("\u23f0 Time Buffer",      "Urgency word frequency",                  "0.44"),
        ("\U0001f60a Satisfaction",  "Positive sentiment ratio",                "0.69"),
        ("\U0001f3af Meaning",      "Long-term goal intents",                  "0.56"),
        ("\U0001f4cd Location",     "Geo-related queries",                     "0.78"),
    ]
    # Table header
    col_h = ["Dimension", "Signal Source", "Value"]
    col_w = [Inches(1.9), Inches(3.0), Inches(0.9)]
    x0 = Inches(0.5)
    x = x0
    for c, w in zip(col_h, col_w):
        add_rect(slide, x, Inches(1.5), w, Inches(0.35), PRIMARY)
        add_text_box(slide, x, Inches(1.52), w, Inches(0.3), c, 9, WHITE, True, PP_ALIGN.CENTER)
        x += w
    # Table rows
    for j, (dim, signal, val) in enumerate(dims):
        top = Inches(1.88) + j * Inches(0.35)
        bg = BG_CARD if j % 2 == 0 else RGBColor(0x15, 0x15, 0x28)
        x = x0
        for i, (v, w) in enumerate(zip([dim, signal, val], col_w)):
            add_rect(slide, x, top, w, Inches(0.33), bg)
            color = WHITE if i == 0 else (GOLD if i == 2 else LIGHT_GRAY)
            add_text_box(slide, x + Inches(0.08), top + Inches(0.02), w - Inches(0.16), Inches(0.28),
                         v, 8, color, i == 0, PP_ALIGN.LEFT if i < 2 else PP_ALIGN.CENTER)
            x += w

    # --- Right panel: How it works ---
    rx = Inches(6.6)
    rw = Inches(6.2)

    # Block 1: Quantification Formula
    add_rect(slide, rx, Inches(1.5), rw, Inches(1.5), BG_CARD, RGBColor(0x33, 0x33, 0x55))
    add_text_box(slide, rx + Inches(0.2), Inches(1.55), rw, Inches(0.3),
                 "Quantification Formula", 13, ACCENT, True)
    add_text_box(slide, rx + Inches(0.2), Inches(1.9), rw - Inches(0.4), Inches(0.3),
                 "x_t = clip( \u03b1 \u00b7 x_{t-1} + (1-\u03b1) \u00b7 signal_t , 0, 1 )", 12, GOLD, True)
    add_text_box(slide, rx + Inches(0.2), Inches(2.25), rw - Inches(0.4), Inches(0.6),
                 "\u03b1 = memory decay rate (per dimension)\n"
                 "signal_t = fused passive input (typing speed, intent keywords, timing)\n"
                 "All inference from keyboard input — zero explicit questionnaires",
                 9, LIGHT_GRAY)

    # Block 2: Particle Filter
    add_rect(slide, rx, Inches(3.15), rw, Inches(1.4), BG_CARD, RGBColor(0x33, 0x33, 0x55))
    add_text_box(slide, rx + Inches(0.2), Inches(3.2), rw, Inches(0.3),
                 "Bayesian Particle Filter (500 particles)", 13, ACCENT, True)
    pf_items = [
        "Each particle = one hypothesis of user's 13-dim state",
        "weight_i \u221d P(observation | state_i) \u00d7 prior_weight_i",
        "Output: weighted mean (point estimate) + variance (uncertainty)",
        "No ground truth needed — converges from typing patterns alone",
    ]
    for i, item in enumerate(pf_items):
        add_text_box(slide, rx + Inches(0.3), Inches(3.55) + i * Inches(0.24), rw - Inches(0.5), Inches(0.22),
                     f"\u2022  {item}", 9, LIGHT_GRAY)

    # Block 3: LIX Bridge
    add_rect(slide, rx, Inches(4.7), rw, Inches(1.6), BG_CARD, GOLD)
    add_text_box(slide, rx + Inches(0.2), Inches(4.75), rw, Inches(0.3),
                 "\U0001f310 Twin \u2192 LIX Marketplace Bridge", 13, GOLD, True)
    bridge_items = [
        ("preferred_domains",   "State dims \u2192 top-3 active life areas"),
        ("risk_tolerance",      "Stress \u00d7 Optionality \u2192 agent Bond threshold"),
        ("price_vs_quality",    "Wealth \u00d7 Time Buffer \u2192 ranking weight flip"),
        ("privacy_mode",        "Only anonymized snapshot sent to agents"),
    ]
    for i, (k, v) in enumerate(bridge_items):
        add_text_box(slide, rx + Inches(0.3), Inches(5.1) + i * Inches(0.28), Inches(2), Inches(0.25),
                     k, 9, WHITE, True)
        add_text_box(slide, rx + Inches(2.5), Inches(5.1) + i * Inches(0.28), Inches(3.5), Inches(0.25),
                     v, 9, LIGHT_GRAY)

    # Bottom: Bellman equation
    add_rect(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.85), RGBColor(0x1A, 0x1A, 0x3E), ACCENT)
    add_text_box(slide, Inches(0.7), Inches(6.55), Inches(4), Inches(0.3),
                 "Bellman Life Optimiser (quasi-hyperbolic):", 11, ACCENT, True)
    add_text_box(slide, Inches(0.7), Inches(6.85), Inches(11.5), Inches(0.4),
                 "Q(s,a) = h(past) + R(s,a) + \u03b2\u00b7\u03b4\u00b7\u03a3 P(s'|s,a)\u00b7V(s')   "
                 "where h(past) = \u03a3 \u03b1^{\u0394t} \u00b7 R_{t-i}  |  "
                 "\u03b2 = present bias  |  \u03b4 = long-run discount  |  \u03b1 = habit decay",
                 10, LIGHT_GRAY)


def build_tech_lix(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "LIX — Lumi Intent Exchange", 30, WHITE, True)
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.6),
                 "Demand-driven marketplace: Buyers set budgets, AI agents broadcast orders, suppliers accept or reject",
                 16, LIGHT_GRAY)

    # Flow
    flow_items = ["Buyer Budget", "→", "Agent Broadcast", "→", "Supplier Match", "→", "Accept/Reject", "→", "Settle"]
    x = Inches(0.4)
    for item in flow_items:
        if item == "→":
            add_text_box(slide, x, Inches(2.0), Inches(0.4), Inches(0.5), "→", 20, ACCENT, True, PP_ALIGN.CENTER)
            x += Inches(0.4)
        else:
            w = Inches(1.8) if len(item) > 12 else Inches(1.5)
            add_rect(slide, x, Inches(1.9), w, Inches(0.7), BG_CARD, PRIMARY)
            add_text_box(slide, x, Inches(2.0), w, Inches(0.5), item, 11, WHITE, False, PP_ALIGN.CENTER)
            x += w + Inches(0.1)

    # Key metrics
    metrics = [
        ("Demand → Match", "~2.5s", "Latency"),
        ("Asset Types", "2+", "Physical & Virtual"),
        ("Accept Fee", "1–3%", "Per Transaction"),
        ("Suppliers", "Anyone", "Open Marketplace"),
    ]
    for i, (label, value, sub) in enumerate(metrics):
        left = Inches(0.8) + i * Inches(3.1)
        card = add_rect(slide, left, Inches(3.2), Inches(2.8), Inches(1.6), BG_CARD, RGBColor(0x33, 0x33, 0x55))
        add_text_box(slide, left, Inches(3.35), Inches(2.8), Inches(0.3), label, 12, MUTED, False, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(3.7), Inches(2.8), Inches(0.6), value, 32, GOLD, True, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(4.3), Inches(2.8), Inches(0.3), sub, 11, MUTED, False, PP_ALIGN.CENTER)

    # Validation stages
    add_text_box(slide, Inches(0.8), Inches(5.2), Inches(10), Inches(0.4),
                 "8-Stage Validation Pipeline", 18, WHITE, True)
    stages = ["Source Allowlist", "Inventory Check", "Expiry Validation", "Trust Score",
              "Price Drift ≤30%", "Budget Overrun", "Delivery Feasibility", "Fraud Detection"]
    for i, s in enumerate(stages):
        left = Inches(0.8) + i * Inches(1.5)
        add_rect(slide, left, Inches(5.7), Inches(1.35), Inches(0.7), BG_CARD, ACCENT)
        add_text_box(slide, left, Inches(5.75), Inches(1.35), Inches(0.6), f"{i+1}. {s}", 9, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
                 "Open to all suppliers: Physical Goods + Virtual Skills (consulting, design) + Intentions (career, investment)",
                 11, MUTED)


# === PART 2: INDUSTRY ===
def build_industry_section(prs):
    add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                       "02", "INDUSTRY",
                       "The Fat-Tail Economy: Scarce AI Capabilities Before AGI")

def build_industry_landscape(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Industry Landscape: The Pre-AGI Window", 30, WHITE, True)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
                 "In LIX, the DEMAND side holds pricing power: buyers set budgets, AI agents broadcast orders, "
                 "and suppliers compete to accept — fundamentally reversing traditional marketplace dynamics.",
                 16, LIGHT_GRAY)

    # Comparison table header
    headers = ["Dimension", "Traditional Long-Tail", "AI Fat-Tail (Intent Market)"]
    for i, h in enumerate(headers):
        left = Inches(0.8) + i * Inches(4)
        w = Inches(3.7) if i > 0 else Inches(3.5)
        add_rect(slide, left, Inches(2.3), w, Inches(0.5), PRIMARY)
        add_text_box(slide, left, Inches(2.33), w, Inches(0.45), h, 12, WHITE, True, PP_ALIGN.CENTER)

    rows = [
        ("Pricing Power", "Supplier sets price", "Demand sets budget, supplier accepts/rejects"),
        ("Order Flow", "Buyer searches → picks", "Buyer broadcasts intent → suppliers compete"),
        ("Matching", "Manual comparison", "AI agent auto-matches best offers to budget"),
        ("Scarcity", "Low (replicable goods)", "High (skills, services, intentions)"),
        ("Network Effect", "Weak bilateral", "Strong trilateral (User-Agent-Supplier)"),
        ("Revenue Concentration", "20% tail → 20% revenue", "20% tail → 80% revenue"),
    ]
    for j, (dim, trad, fat) in enumerate(rows):
        top = Inches(2.85) + j * Inches(0.58)
        bg = BG_CARD if j % 2 == 0 else RGBColor(0x15, 0x15, 0x28)
        for i, val in enumerate([dim, trad, fat]):
            left = Inches(0.8) + i * Inches(4)
            w = Inches(3.7) if i > 0 else Inches(3.5)
            add_rect(slide, left, top, w, Inches(0.55), bg)
            color = WHITE if i == 0 else (GOLD if i == 2 else LIGHT_GRAY)
            add_text_box(slide, left + Inches(0.15), top + Inches(0.05), w - Inches(0.3), Inches(0.45),
                         val, 11, color, i == 2)


def build_industry_demand(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Fat-Tail Demand Composition", 30, WHITE, True)

    categories = [
        ("Professional Consulting", "30%", "$500–5,000/task", "Legal, medical, tax planning", PRIMARY),
        ("Creative Production", "25%", "$200–2,000/task", "Design, writing, architecture", ACCENT),
        ("Complex Decisions", "20%", "$1,000–10,000/task", "Investment, career, education", GREEN),
        ("Personalized Services", "15%", "$100–1,000/task", "Fitness, therapy, coaching", GOLD),
        ("Urgent Tasks", "10%", "2–5× market rate", "Emergency translation, crisis PR", PINK),
    ]
    for i, (name, share, price, examples, color) in enumerate(categories):
        top = Inches(1.3) + i * Inches(1.15)
        # Percentage bar background
        add_rect(slide, Inches(0.8), top, Inches(12), Inches(1.0), BG_CARD, RGBColor(0x33, 0x33, 0x55))
        # Color bar proportional
        pct = int(share.replace('%', ''))
        bar_w = Inches(12 * pct / 100)
        r, g, b = int(str(color)[0:2], 16), int(str(color)[2:4], 16), int(str(color)[4:6], 16)
        add_rect(slide, Inches(0.8), top, bar_w, Inches(1.0), RGBColor(r // 3, g // 3, b // 3))
        add_rect(slide, Inches(0.8), top, Inches(0.1), Inches(1.0), color)
        add_text_box(slide, Inches(1.1), top + Inches(0.08), Inches(3.5), Inches(0.4), name, 16, color, True)
        add_text_box(slide, Inches(5), top + Inches(0.08), Inches(1.5), Inches(0.4), share, 22, WHITE, True)
        add_text_box(slide, Inches(7), top + Inches(0.08), Inches(2.5), Inches(0.4), price, 14, GOLD, True)
        add_text_box(slide, Inches(1.1), top + Inches(0.55), Inches(10), Inches(0.35), examples, 12, MUTED)

    add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.4),
                 "Key Insight: These demands cannot be perfectly solved by standardized AI before AGI — "
                 "requiring specialized human agents or vertical AI.", 13, GOLD, True, PP_ALIGN.CENTER)


# === PART 3: MARKET ===
def build_market_section(prs):
    add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                       "03", "MARKET",
                       "Total Addressable Market × Market Share × Growth Trajectory")

def build_market_tam(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Total Addressable Market: $1.64 Trillion", 30, WHITE, True)

    markets = [
        ("Financial Services", "$20T", "25%", "15%", "$750B"),
        ("Education & Training", "$5T", "40%", "20%", "$400B"),
        ("Creative & Design", "$2T", "50%", "25%", "$250B"),
        ("Healthcare", "$8T", "20%", "10%", "$160B"),
        ("Legal Services", "$1T", "30%", "20%", "$60B"),
        ("Enterprise Consulting", "$300B", "35%", "20%", "$21B"),
    ]
    # Header
    cols = ["Sector", "Global Size", "AI-able %", "Capturable %", "Lumi TAM"]
    widths = [Inches(3), Inches(2), Inches(2), Inches(2.5), Inches(2.5)]
    x = Inches(0.6)
    for c, w in zip(cols, widths):
        add_rect(slide, x, Inches(1.3), w, Inches(0.5), PRIMARY)
        add_text_box(slide, x, Inches(1.33), w, Inches(0.45), c, 12, WHITE, True, PP_ALIGN.CENTER)
        x += w

    for j, (sector, size, ai_pct, cap_pct, tam) in enumerate(markets):
        top = Inches(1.85) + j * Inches(0.58)
        bg = BG_CARD if j % 2 == 0 else RGBColor(0x15, 0x15, 0x28)
        x = Inches(0.6)
        for i, (val, w) in enumerate(zip([sector, size, ai_pct, cap_pct, tam], widths)):
            add_rect(slide, x, top, w, Inches(0.55), bg)
            color = WHITE if i == 0 else (GOLD if i == 4 else LIGHT_GRAY)
            add_text_box(slide, x + Inches(0.1), top + Inches(0.05), w - Inches(0.2), Inches(0.45),
                         val, 12 if i == 0 else 13, color, i == 4, PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
            x += w

    # Total row
    top = Inches(1.85) + 6 * Inches(0.58)
    x = Inches(0.6)
    for i, (val, w) in enumerate(zip(["TOTAL", "$36.3T", "~30%", "~18%", "$1.64T"], widths)):
        add_rect(slide, x, top, w, Inches(0.6), RGBColor(0x2D, 0x2D, 0x5E))
        add_text_box(slide, x + Inches(0.1), top + Inches(0.08), w - Inches(0.2), Inches(0.45),
                     val, 14, GOLD if i == 4 else WHITE, True, PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
        x += w

def build_market_customers(prs):
    """Target customer segments and why Lumi can capture this market."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Target Customers & Why We Win", 30, WHITE, True)
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
                 "Who uses Lumi — and why they never switch once they start",
                 15, MUTED)

    # 5 customer personas
    personas = [
        ("🎓 Students & Researchers", "18–30", "High",
         "Decision-dense academic life: course selection, career planning, grad school applications"),
        ("💼 Knowledge Workers", "25–45", "Very High",
         "Professionals making 50+ decisions/day: scheduling, sourcing, comparing services"),
        ("🚀 Freelancers & Creators", "22–40", "High",
         "Pricing, client matching, project scoping — their income depends on smart intent execution"),
        ("🏢 Small Business Owners", "30–55", "Very High",
         "Procurement, vendor selection, budget optimization — every dollar decision matters"),
        ("🛒 Smart Consumers", "20–50", "Medium-High",
         "High-value purchases (electronics, travel, insurance) where budget-driven matching saves money"),
    ]
    cols_h = ["Segment", "Age", "Intent Density", "Why Lumi Fits"]
    widths = [Inches(2.8), Inches(1), Inches(1.5), Inches(7)]
    x = Inches(0.5)
    for c, w in zip(cols_h, widths):
        add_rect(slide, x, Inches(1.6), w, Inches(0.45), PRIMARY)
        add_text_box(slide, x, Inches(1.62), w, Inches(0.4), c, 11, WHITE, True, PP_ALIGN.CENTER)
        x += w
    for j, (seg, age, density, reason) in enumerate(personas):
        top = Inches(2.1) + j * Inches(0.58)
        bg = BG_CARD if j % 2 == 0 else RGBColor(0x15, 0x15, 0x28)
        x = Inches(0.5)
        for i, (val, w) in enumerate(zip([seg, age, density, reason], widths)):
            add_rect(slide, x, top, w, Inches(0.55), bg)
            color = WHITE if i == 0 else (GOLD if i == 2 else LIGHT_GRAY)
            add_text_box(slide, x + Inches(0.1), top + Inches(0.08), w - Inches(0.2), Inches(0.4),
                         val, 11, color, i == 0, PP_ALIGN.LEFT if i in (0, 3) else PP_ALIGN.CENTER)
            x += w

    # Core characteristics box
    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(6), Inches(0.4),
                 "Collective Customer Profile", 18, WHITE, True)
    traits = [
        "•  Mobile-first: 4+ hrs/day keyboard usage",
        "•  Decision-dense: 30–80 intent-bearing actions/day",
        "•  Price-sensitive but time-poor: willing to pay for better outcomes",
        "•  Privacy-conscious: prefer local-first AI over cloud data sharing",
    ]
    for i, t in enumerate(traits):
        add_text_box(slide, Inches(0.7), Inches(5.5) + i * Inches(0.35), Inches(5.5), Inches(0.3),
                     t, 11, LIGHT_GRAY)

    # Why we win box
    add_rect(slide, Inches(6.8), Inches(5.1), Inches(5.8), Inches(2.2), BG_CARD, GOLD)
    add_text_box(slide, Inches(7.1), Inches(5.2), Inches(5.2), Inches(0.4),
                 "🏆 Why Lumi Captures This Market", 16, GOLD, True)
    reasons = [
        "•  Zero-friction entry: keyboard = universal input layer",
        "•  Digital Twin lock-in: switching cost grows over time",
        "•  Budget-driven matching: users save money instantly",
        "•  Network effect: more users = better supplier competition",
        "•  Habit formation: daily keyboard use = daily engagement",
    ]
    for i, r in enumerate(reasons):
        add_text_box(slide, Inches(7.3), Inches(5.65) + i * Inches(0.32), Inches(5.1), Inches(0.3),
                     r, 11, LIGHT_GRAY)


def build_market_why_now(prs):
    """Why Now — OpenClaw has validated PMF for the new work logic."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "WHY NOW", 14, ACCENT, True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(1),
                 "OpenClaw Has Already Proven PMF — Go Straight to Growth", 32, WHITE, True)

    # Left card: OpenClaw PMF Evidence
    add_rect(slide, Inches(0.6), Inches(1.9), Inches(5.8), Inches(4.8), BG_CARD, GREEN)
    add_text_box(slide, Inches(0.9), Inches(2.05), Inches(5.2), Inches(0.5),
                 "✅ OpenClaw: PMF Validated", 20, GREEN, True)
    add_text_box(slide, Inches(0.9), Inches(2.6), Inches(5.2), Inches(0.5),
                 "OpenClaw proved that the new AI-Agent work logic works:", 13, LIGHT_GRAY)

    evidence = [
        ("New Work Logic", "Intent-Driven → Agent Execution → Supply-Demand Match"),
        ("User Behavior", "Users naturally adopt agent-based workflows over manual search"),
        ("Retention Signal", "High repeat usage — users return because the logic saves time & money"),
        ("Market Response", "Organic demand growth without paid acquisition — pull, not push"),
    ]
    for i, (label, desc) in enumerate(evidence):
        top = Inches(3.2) + i * Inches(0.75)
        add_rect(slide, Inches(0.9), top, Inches(0.1), Inches(0.6), GREEN)
        add_text_box(slide, Inches(1.2), top + Inches(0.02), Inches(4.8), Inches(0.3),
                     label, 14, GOLD, True)
        add_text_box(slide, Inches(1.2), top + Inches(0.32), Inches(4.8), Inches(0.3),
                     desc, 11, LIGHT_GRAY)

    # Core conclusion
    add_rect(slide, Inches(0.9), Inches(6.25), Inches(5.2), Inches(0.35), RGBColor(0x10, 0x3D, 0x2D))
    add_text_box(slide, Inches(1.0), Inches(6.27), Inches(5.0), Inches(0.3),
                 "\"Market exists. Logic works. Users stay.\"", 13, GREEN, True, PP_ALIGN.CENTER)

    # Right card: Why Growth Phase Now
    add_rect(slide, Inches(6.8), Inches(1.9), Inches(5.8), Inches(4.8), BG_CARD, GOLD)
    add_text_box(slide, Inches(7.1), Inches(2.05), Inches(5.2), Inches(0.5),
                 "🚀 Skip Validation → Enter Growth", 20, GOLD, True)
    add_text_box(slide, Inches(7.1), Inches(2.6), Inches(5.2), Inches(0.5),
                 "Why Lumi doesn't need to prove the market — it's already proven:", 13, LIGHT_GRAY)

    reasons = [
        ("✅ PMF Validated", "OpenClaw proved intent-driven agent workflows have real demand — no exploration needed"),
        ("✅ User Behavior Formed", "The new work logic (AI does the work, user sets the goal) is accepted and habitual"),
        ("✅ Perfect Timing", "AI Agent wave is peaking — the window to capture mind-share is NOW, not later"),
        ("✅ Direct to Scale", "With proven logic, every dollar goes to growth, not experimentation"),
    ]
    for i, (title, desc) in enumerate(reasons):
        top = Inches(3.2) + i * Inches(0.75)
        add_text_box(slide, Inches(7.3), top + Inches(0.02), Inches(5.0), Inches(0.3),
                     title, 14, WHITE, True)
        add_text_box(slide, Inches(7.3), top + Inches(0.32), Inches(5.0), Inches(0.3),
                     desc, 11, LIGHT_GRAY)

    # Bottom banner
    add_rect(slide, Inches(0.6), Inches(6.9), Inches(12), Inches(0.45), RGBColor(0x2D, 0x2D, 0x5E))
    add_text_box(slide, Inches(0.6), Inches(6.93), Inches(12), Inches(0.4),
                 "OpenClaw proved the logic.  Lumi scales the network.  The market share is ours to take.",
                 15, GOLD, True, PP_ALIGN.CENTER)


def build_market_share(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Market Share Trajectory (Mind-Share Driven Growth)", 28, WHITE, True)

    # Growth data
    years = [
        ("Y1", "5B", "2%", "1%", "1M", "Early adopters, developer community"),
        ("Y2", "5.2B", "5%", "2%", "5M", "Productivity bloggers, word-of-mouth"),
        ("Y3*", "5.4B", "10%", "4%", "20M", "TIPPING POINT — mass adoption begins"),
        ("Y4", "5.6B", "15%", "6%", "50M", "Mainstream growth"),
        ("Y5", "5.8B", "20%", "8%", "100M", "Platform-level scale"),
        ("Y7", "6.2B", "30%", "10%", "200M", "Intent economy infrastructure"),
    ]
    cols_h = ["Year", "Global Users", "Agent OS %", "Lumi Share", "Lumi Users", "Phase"]
    widths = [Inches(1), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(4)]
    x = Inches(0.5)
    for c, w in zip(cols_h, widths):
        add_rect(slide, x, Inches(1.3), w, Inches(0.5), PRIMARY)
        add_text_box(slide, x, Inches(1.33), w, Inches(0.45), c, 11, WHITE, True, PP_ALIGN.CENTER)
        x += w

    for j, row in enumerate(years):
        top = Inches(1.85) + j * Inches(0.65)
        bg = RGBColor(0x1E, 0x1E, 0x40) if row[0] in ("Y3*",) else (BG_CARD if j % 2 == 0 else RGBColor(0x15, 0x15, 0x28))
        x = Inches(0.5)
        for i, (val, w) in enumerate(zip(row, widths)):
            add_rect(slide, x, top, w, Inches(0.6), bg)
            is_highlight = row[0] == "Y3*"
            color = GOLD if is_highlight else (WHITE if i in (0, 4) else LIGHT_GRAY)
            add_text_box(slide, x + Inches(0.05), top + Inches(0.08), w - Inches(0.1), Inches(0.45),
                         val, 11, color, is_highlight or i == 4, PP_ALIGN.CENTER if i < 5 else PP_ALIGN.LEFT)
            x += w

    add_text_box(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(0.5),
                 "* Y3 = Tipping Point — comparable to TikTok's 2019 inflection, mass adoption begins",
                 14, GOLD, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.8),
                 "Defensibility: Digital Twin lock-in (switching cost) + Habit path dependency + "
                 "Agent marketplace ecosystem + First-mover brand (\"Agent = Lumi\")",
                 13, MUTED, False, PP_ALIGN.CENTER)


# === PART 4: FINANCING ===
def build_financing_section(prs):
    add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                       "04", "FINANCING PLAN",
                       "Seed Round: $1,000,000 for 10% Equity  |  Valuation: $10,000,000")

def build_business_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Business Model: Dual Revenue Engine", 30, WHITE, True)

    # Revenue stream 1
    add_rect(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(2.8), BG_CARD, PRIMARY)
    add_text_box(slide, Inches(0.9), Inches(1.45), Inches(5), Inches(0.5),
                 "💎 Stream 1: Keyboard Subscription", 20, PRIMARY, True)
    items1 = [
        "Freemium model: Free basic → Premium AI features",
        "Premium: $9.99–19.99/month per user",
        "Features: Digital Twin, Destiny Engine, priority Agent",
        "Target: 20–40% conversion rate at scale",
        "Recurring SaaS revenue with high retention",
    ]
    for i, item in enumerate(items1):
        add_text_box(slide, Inches(1.1), Inches(2.1) + i * Inches(0.38), Inches(5), Inches(0.35),
                     f"•  {item}", 12, LIGHT_GRAY)

    # Revenue stream 2
    add_rect(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.8), BG_CARD, GOLD)
    add_text_box(slide, Inches(7.1), Inches(1.45), Inches(5), Inches(0.5),
                 "🏪 Stream 2: LIX Accept Fee", 20, GOLD, True)
    items2 = [
        "1%–3% transaction fee on matched orders",
        "Anyone can be a supplier — open marketplace",
        "Revenue scales with GMV (network effect)",
        "8-stage validation ensures quality → trust → volume",
        "Goods + skills + services + intentions",
    ]
    for i, item in enumerate(items2):
        add_text_box(slide, Inches(7.3), Inches(2.1) + i * Inches(0.38), Inches(5), Inches(0.35),
                     f"•  {item}", 12, LIGHT_GRAY)

    # Revenue projection
    add_text_box(slide, Inches(0.8), Inches(4.4), Inches(10), Inches(0.5),
                 "Revenue Projection", 22, WHITE, True)
    rev_cols = ["Year", "Subscription", "LIX GMV", "Platform Fee (1-3%)", "Total Revenue"]
    rev_widths = [Inches(1.5), Inches(2.5), Inches(2.5), Inches(3), Inches(2.5)]
    x = Inches(0.6)
    for c, w in zip(rev_cols, rev_widths):
        add_rect(slide, x, Inches(4.9), w, Inches(0.45), PRIMARY)
        add_text_box(slide, x, Inches(4.92), w, Inches(0.4), c, 11, WHITE, True, PP_ALIGN.CENTER)
        x += w
    rev_rows = [
        ("Y1", "$12M", "$2M", "$60K", "$12M"),
        ("Y3", "$300M", "$500M", "$15M", "$315M"),
        ("Y5", "$2B", "$10B", "$300M", "$2.3B"),
    ]
    for j, row in enumerate(rev_rows):
        top = Inches(5.4) + j * Inches(0.5)
        bg = BG_CARD if j % 2 == 0 else RGBColor(0x15, 0x15, 0x28)
        x = Inches(0.6)
        for i, (v, w) in enumerate(zip(row, rev_widths)):
            add_rect(slide, x, top, w, Inches(0.48), bg)
            color = GOLD if i == 4 else (WHITE if i == 0 else LIGHT_GRAY)
            add_text_box(slide, x, top + Inches(0.03), w, Inches(0.4), v, 12, color, i == 4, PP_ALIGN.CENTER)
            x += w

    add_text_box(slide, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
                 "Key Inflection: By Y5, LIX marketplace surpasses subscription revenue — platform economics at scale",
                 12, GOLD, True, PP_ALIGN.CENTER)


def build_use_of_funds(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Use of Funds & Investor Returns", 30, WHITE, True)

    # Deal terms
    add_rect(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(2.5), BG_CARD, PRIMARY)
    add_text_box(slide, Inches(0.9), Inches(1.45), Inches(5), Inches(0.5),
                 "Seed Round Terms", 20, PRIMARY, True)
    terms = [
        ("Raising", "$1,000,000"),
        ("Equity", "10%"),
        ("Pre-Money Valuation", "$10,000,000"),
        ("Post-Money Valuation", "$11,000,000"),
    ]
    for i, (label, val) in enumerate(terms):
        add_text_box(slide, Inches(1.1), Inches(2.1) + i * Inches(0.45), Inches(3), Inches(0.4),
                     label, 14, LIGHT_GRAY)
        add_text_box(slide, Inches(4.2), Inches(2.1) + i * Inches(0.45), Inches(2), Inches(0.4),
                     val, 14, GOLD, True)

    # Use of funds — detailed breakdown table
    add_rect(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.8), BG_CARD, ACCENT)
    add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5), Inches(0.4),
                 "Fund Allocation — Detailed Breakdown", 18, ACCENT, True)
    allocs = [
        ("Office Space", "$50,000", "5%", LIGHT_GRAY),
        ("Equipment & Infrastructure", "$100,000", "10%", LIGHT_GRAY),
        ("Edge Model Training", "$100,000", "10%", PRIMARY),
        ("CTO / Tech Lead", "$100,000", "10%", PRIMARY),
        ("Legal & Compliance", "$100,000", "10%", LIGHT_GRAY),
        ("Marketing & User Acquisition", "$200,000", "20%", GOLD),
        ("Working Capital (Reserve)", "$350,000", "35%", GREEN),
    ]
    for i, (label, amt, pct, color) in enumerate(allocs):
        top = Inches(1.85) + i * Inches(0.3)
        add_text_box(slide, Inches(7.2), top, Inches(2.8), Inches(0.28), label, 10, color, False)
        add_text_box(slide, Inches(10.2), top, Inches(1.1), Inches(0.28), amt, 10, WHITE, True)
        add_text_box(slide, Inches(11.5), top, Inches(0.8), Inches(0.28), pct, 10, GOLD, True)
    # Total line
    add_rect(slide, Inches(7.1), Inches(3.98), Inches(5.3), Inches(0.02), ACCENT)
    add_text_box(slide, Inches(7.2), Inches(4.0), Inches(2.8), Inches(0.28), "TOTAL", 11, WHITE, True)
    add_text_box(slide, Inches(10.2), Inches(4.0), Inches(1.1), Inches(0.28), "$1,000,000", 11, GOLD, True)
    add_text_box(slide, Inches(11.5), Inches(4.0), Inches(0.8), Inches(0.28), "100%", 11, GOLD, True)

    # ROI
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.5),
                 "Investor Return Scenarios", 22, WHITE, True)
    scenarios = [
        ("Conservative", "30%", "Niche tool (<50M users)", "3×", "10×", LIGHT_GRAY),
        ("Base Case", "45%", "Mainstream platform (100–500M)", "10×", "50×", WHITE),
        ("Bull Case", "25%", "Infrastructure (500M+ users)", "50×", "200×", GOLD),
    ]
    cols_s = ["Scenario", "Probability", "Outcome", "Return", "At Scale", ""]
    ws = [Inches(1.8), Inches(1.5), Inches(4), Inches(1.5), Inches(1.5)]
    x = Inches(0.8)
    for c, w in zip(cols_s[:5], ws):
        add_rect(slide, x, Inches(4.8), w, Inches(0.45), PRIMARY)
        add_text_box(slide, x, Inches(4.82), w, Inches(0.4), c, 11, WHITE, True, PP_ALIGN.CENTER)
        x += w
    for j, (scen, prob, outcome, ret, scale, color) in enumerate(scenarios):
        top = Inches(5.3) + j * Inches(0.55)
        bg = BG_CARD if j % 2 == 0 else RGBColor(0x15, 0x15, 0x28)
        x = Inches(0.8)
        for i, (v, w) in enumerate(zip([scen, prob, outcome, ret, scale], ws)):
            add_rect(slide, x, top, w, Inches(0.5), bg)
            c = GOLD if i >= 3 else (color if i == 0 else LIGHT_GRAY)
            add_text_box(slide, x + Inches(0.05), top + Inches(0.05), w - Inches(0.1), Inches(0.4),
                         v, 11, c, i >= 3, PP_ALIGN.CENTER if i != 2 else PP_ALIGN.LEFT)
            x += w

    add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.4),
                 "Expected Return (Weighted): 15.4× — based on risk-adjusted probability analysis",
                 14, GOLD, True, PP_ALIGN.CENTER)


def build_milestones(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Key Milestones & Enterprise Pipeline", 30, WHITE, True)

    milestones = [
        ("2024", "Founded in London, UK", "Company incorporation, initial R&D"),
        ("Dec 2025", "Demo product concept", "Web simulator + Android keyboard prototype"),
        ("Feb 2026", "Product demo ready", "Full LIX MVP, 4-layer agent, DTOE v0.3"),
        ("Q2 2026", "Seed round close", "1M users target, LIX marketplace live"),
        ("Q4 2026", "Growth phase", "5M users, expand supply network"),
        ("2027", "Series A ready", "20M users, tipping point"),
    ]
    for i, (date, event, desc) in enumerate(milestones):
        top = Inches(1.3) + i * Inches(0.95)
        # Timeline line
        add_rect(slide, Inches(2.5), top, Inches(0.04), Inches(0.95), RGBColor(0x33, 0x33, 0x55))
        # Dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.38), top + Inches(0.2), Inches(0.28), Inches(0.28))
        dot.fill.solid()
        dot.fill.fore_color.rgb = PRIMARY
        dot.line.fill.background()
        # Date
        add_text_box(slide, Inches(0.6), top + Inches(0.15), Inches(1.7), Inches(0.4),
                     date, 14, GOLD, True, PP_ALIGN.RIGHT)
        # Event
        add_text_box(slide, Inches(3.0), top + Inches(0.1), Inches(4), Inches(0.35),
                     event, 16, WHITE, True)
        add_text_box(slide, Inches(3.0), top + Inches(0.48), Inches(9), Inches(0.35),
                     desc, 12, MUTED)

    # Enterprise pipeline note
    add_rect(slide, Inches(8), Inches(1.5), Inches(4.5), Inches(2.5), BG_CARD, GOLD)
    add_text_box(slide, Inches(8.3), Inches(1.65), Inches(4), Inches(0.4),
                 "🏦 Enterprise Pipeline", 16, GOLD, True)
    add_text_box(slide, Inches(8.3), Inches(2.1), Inches(4), Inches(1.6),
                 "• Morgan Stanley — Bank-grade AI execution layer\n"
                 "• 90-day PoC framework defined\n"
                 "• Wealth Management & Research workflows\n"
                 "• Digital Twin optimization for advisors",
                 12, LIGHT_GRAY)


# === PART 5: TEAM ===
def build_team_section(prs):
    add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                       "05", "TEAM",
                       "UCL Innovation × Venture Capital Expertise")

def build_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Founding Team", 30, WHITE, True)

    # Songyi Li
    add_rect(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(4.5), BG_CARD, GOLD)
    # Avatar circle
    avatar1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.6), Inches(1.6), Inches(1.2), Inches(1.2))
    avatar1.fill.solid()
    avatar1.fill.fore_color.rgb = GOLD
    avatar1.line.fill.background()
    add_text_box(slide, Inches(1.6), Inches(1.85), Inches(1.2), Inches(0.7),
                 "SL", 28, WHITE, True, PP_ALIGN.CENTER)

    add_text_box(slide, Inches(3.1), Inches(1.6), Inches(3), Inches(0.4), "Songyi Li", 24, WHITE, True)
    add_text_box(slide, Inches(3.1), Inches(2.1), Inches(3), Inches(0.35), "Founder & CEO", 16, GOLD, True)

    edu1 = [
        "🎓 Ph.D. Candidate, FinTech — UCL",
        "🎓 M.Sc. Venture Capital & PE — UCL",
        "🎓 B.Sc. Mathematics — XJTLU & Liverpool",
    ]
    for i, e in enumerate(edu1):
        add_text_box(slide, Inches(1.0), Inches(3.1) + i * Inches(0.4), Inches(5), Inches(0.35),
                     e, 12, LIGHT_GRAY)
    add_text_box(slide, Inches(1.0), Inches(4.5), Inches(5), Inches(0.6),
                 "Leading the convergence of FinTech innovation and AI-powered productivity tools. "
                 "Deep expertise in financial technology research and venture capital.",
                 11, MUTED)

    # Bo Zhao
    add_rect(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(4.5), BG_CARD, ACCENT)
    avatar2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.8), Inches(1.6), Inches(1.2), Inches(1.2))
    avatar2.fill.solid()
    avatar2.fill.fore_color.rgb = ACCENT
    avatar2.line.fill.background()
    add_text_box(slide, Inches(7.8), Inches(1.85), Inches(1.2), Inches(0.7),
                 "BZ", 28, WHITE, True, PP_ALIGN.CENTER)

    add_text_box(slide, Inches(9.3), Inches(1.6), Inches(3), Inches(0.4), "Bo Zhao", 24, WHITE, True)
    add_text_box(slide, Inches(9.3), Inches(2.1), Inches(3), Inches(0.35), "Co-Founder & COO", 16, ACCENT, True)

    edu2 = [
        "🎓 M.Sc. Venture Capital & PE — UCL",
        "🎓 B.Sc. Finance — Durham University",
        "💼 Former financial professional, Guosen Securities",
    ]
    for i, e in enumerate(edu2):
        add_text_box(slide, Inches(7.2), Inches(3.1) + i * Inches(0.4), Inches(5), Inches(0.35),
                     e, 12, LIGHT_GRAY)
    add_text_box(slide, Inches(7.2), Inches(4.5), Inches(5), Inches(0.6),
                 "Bringing venture capital expertise and financial acumen to scale Lumi.AI globally. "
                 "Deep experience in fundraising and financial operations.",
                 11, MUTED)

    # Academic partner
    add_rect(slide, Inches(0.6), Inches(6.1), Inches(12), Inches(1.0), BG_CARD, RGBColor(0x33, 0x33, 0x55))
    add_text_box(slide, Inches(1.0), Inches(6.2), Inches(5), Inches(0.4),
                 "🏛️  Academic Partner: University College London (UCL)", 16, WHITE, True)
    add_text_box(slide, Inches(1.0), Inches(6.6), Inches(11), Inches(0.35),
                 "Top 10 global university  |  FinTech research hub  |  London, UK  |  Focus: FinTech × AI × Privacy Tech",
                 12, MUTED)


def build_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), ACCENT)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                 "🌟 Lumi.AI", 36, PRIMARY, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
                 "Building the Infrastructure of the Intent Economy", 32, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1),
                 "\"Lumi is not selling a tool — it's building the network that distributes "
                 "scarce AI capabilities before AGI arrives.\"",
                 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # Key numbers
    nums = [
        ("$10M", "Valuation"),
        ("$1M", "Raising"),
        ("10%", "Equity"),
        ("15.4×", "Expected Return"),
    ]
    for i, (num, label) in enumerate(nums):
        left = Inches(1.5) + i * Inches(2.8)
        add_rect(slide, left, Inches(5.0), Inches(2.3), Inches(1.2), BG_CARD, PRIMARY)
        add_text_box(slide, left, Inches(5.1), Inches(2.3), Inches(0.6),
                     num, 28, GOLD, True, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(5.7), Inches(2.3), Inches(0.4),
                     label, 13, MUTED, False, PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(6.6), Inches(9), Inches(0.5),
                 "Contact: team@lumi.ai  |  lumi-agent-simulator.vercel.app  |  London, UK",
                 14, MUTED, False, PP_ALIGN.CENTER)


# ========== MAIN ==========
def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # Cover + Problem + Solution
    build_cover(prs)
    build_problem(prs)
    build_solution(prs)

    # Part 1: Technology
    build_tech_section(prs)
    build_tech_arch_diagram(prs)
    build_tech_arch(prs)
    build_tech_digital_twin(prs)
    build_tech_lix(prs)

    # Part 2: Industry
    build_industry_section(prs)
    build_industry_landscape(prs)
    build_industry_demand(prs)

    # Part 3: Market
    build_market_section(prs)
    build_market_tam(prs)
    build_market_customers(prs)
    build_market_why_now(prs)
    build_market_share(prs)

    # Part 4: Financing Plan
    build_financing_section(prs)
    build_business_model(prs)
    build_use_of_funds(prs)

    # Part 5: Team
    build_team_section(prs)
    build_team(prs)

    # Closing
    build_closing(prs)

    output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)),
                               "Lumi_AI_Investor_Pitch_Deck.pptx")
    prs.save(output_path)
    print(f"✅ Pitch deck saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
